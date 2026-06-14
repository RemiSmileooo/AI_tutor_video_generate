"""导出可编辑的 PowerPoint(.pptx)。

依据课程结构生成标准 16:9 PPT：
- 封面页：居中大标题 + 副标题
- 内容页：标题 + 强调短线 + 要点列表
- 每页的口播稿写入"演讲者备注"，方便老师拿来直接讲

配色跟随当前主题(THEME)，与视频中的 PPT 风格保持一致。
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .config import settings, THEMES
from .models import Course


def _rgb(t: tuple) -> RGBColor:
    return RGBColor(t[0], t[1], t[2])


def _set_bg(slide, color: tuple) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _no_line(shape) -> None:
    shape.line.fill.background()


def export_pptx(course: Course, out_path: str | Path, theme_name: str | None = None) -> str:
    th = THEMES.get(theme_name or settings.theme_name, THEMES["apple"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    for slide_data in course.slides:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, th.bg_top)

        if slide_data.kind == "cover":
            _build_cover(slide, slide_data, th, W, H)
        else:
            _build_content(slide, slide_data, th, W, H)

        # 演讲者备注 = 该页所有口播稿
        scripts = [seg.script for seg in slide_data.segments if seg.script]
        if scripts:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(scripts)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return str(out_path)


def _build_cover(slide, data, th, W, H) -> None:
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), W - Inches(2.0), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = _rgb(th.title)

    # 强调短线
    line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, W / 2 - Inches(0.9), Inches(4.45), Inches(1.8), Pt(6))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(th.accent)
    _no_line(line)

    subtitle = (data.bullets[0] if data.bullets else "") or ""
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), W - Inches(2.0), Inches(0.8))
        sp = sb.text_frame.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        r = sp.add_run()
        r.text = subtitle
        r.font.size = Pt(20)
        r.font.color.rgb = _rgb(th.text)


def _build_content(slide, data, th, W, H) -> None:
    # 标题
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), W - Inches(1.8), Inches(1.0))
    tp = tb.text_frame.paragraphs[0]
    r = tp.add_run()
    r.text = data.title
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = _rgb(th.title)

    # 标题下强调短线
    line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.55), Inches(2.2), Pt(6))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(th.accent)
    _no_line(line)

    if not data.bullets:
        return

    # 要点列表
    body = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), W - Inches(2.2), H - Inches(2.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(data.bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        # 圆点
        dot = p.add_run()
        dot.text = "●  "
        dot.font.size = Pt(16)
        dot.font.color.rgb = _rgb(th.accent)
        # 文本
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(22)
        run.font.color.rgb = _rgb(th.text)
